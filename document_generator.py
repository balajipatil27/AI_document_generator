from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class DocumentGenerator:
    @staticmethod
    def generate_docx(project_data, contents):
        doc = Document()
        
        # Set document styles
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Calibri'
        font.size = Pt(11)
        
        # Title
        title = doc.add_heading(project_data['topic'], 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add content for each section
        for section in project_data.get('outline', []):
            section_id = section.get('id')
            section_title = section.get('title')
            
            content = next((c for c in contents if c['section_id'] == section_id), None)
            
            if content:
                # Add section heading
                heading = doc.add_heading(section_title, level=1)
                
                # Add content with proper formatting
                paragraphs = content['content_text'].split('\n')
                for paragraph in paragraphs:
                    if paragraph.strip():
                        p = doc.add_paragraph(paragraph.strip())
                        p.style = doc.styles['Normal']
        
        return doc
    
    @staticmethod
    def generate_pptx(project_data, contents):
        prs = Presentation()
        
        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = project_data['topic']
        subtitle_shape.text = "AI-Generated Presentation"
        
        # Content slides
        for i, section in enumerate(project_data.get('outline', [])):
            section_id = section.get('id')
            section_title = section.get('title')
            
            content = next((c for c in contents if c['section_id'] == section_id), None)
            
            if content and i > 0:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                
                # Title
                title_shape = slide.shapes.title
                title_shape.text = section_title
                
                # Content
                content_shape = slide.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.text = content['content_text']
                
                # Format text
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(18)
                    paragraph.font.name = 'Calibri'
        
        return prs